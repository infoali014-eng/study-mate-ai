import base64
import html
from pathlib import Path

import streamlit as st


def file_exists(file_path):
    """Return True when the stored original file path still exists."""
    return bool(file_path) and Path(file_path).is_file()


def get_file_bytes(file_path):
    """Read a file as bytes for preview or download."""
    return Path(file_path).read_bytes()


def get_file_download_button(file_path, label="Download Original File"):
    """Render a Streamlit download button for any stored file."""
    path = Path(file_path)
    if not path.is_file():
        st.error("Original file not found. Please re-upload this document.")
        return

    st.download_button(
        label=label,
        data=get_file_bytes(path),
        file_name=path.name,
        mime="application/octet-stream",
        use_container_width=True,
    )


def preview_pdf(file_path, height=720):
    """Show a PDF inside the Streamlit app with safe fallbacks."""
    pdf_bytes = get_file_bytes(file_path)

    try:
        st.pdf(pdf_bytes, height=height)
        return
    except Exception:
        # Older deployments may not have Streamlit's PDF extra installed.
        # The base64 iframe below keeps preview working without crashing.
        pass

    if len(pdf_bytes) > 25 * 1024 * 1024:
        st.warning(
            "This PDF is large, so embedded preview is unavailable in this environment. "
            "Use the download button below to open the original file."
        )
        return

    pdf_base64 = base64.b64encode(pdf_bytes).decode("utf-8")
    st.markdown(
        f"""
        <div class="preview-frame">
            <iframe
                src="data:application/pdf;base64,{pdf_base64}"
                width="100%"
                height="{height}"
                type="application/pdf">
            </iframe>
        </div>
        """,
        unsafe_allow_html=True,
    )


def preview_text_file(file_path, max_characters=20000):
    """Show a readable scrollable preview for TXT-like files."""
    text = Path(file_path).read_text(encoding="utf-8", errors="ignore")
    preview_text = html.escape(text[:max_characters])

    st.markdown(
        f"""
        <div class="text-preview">
            <pre>{preview_text}</pre>
        </div>
        """,
        unsafe_allow_html=True,
    )


def preview_image(file_path):
    """Display an uploaded image inside the app."""
    st.image(str(file_path), use_container_width=True)


def preview_extracted_text(text, max_characters=20000):
    """Show extracted text in the same scrollable preview style."""
    preview_text = html.escape((text or "")[:max_characters])
    st.markdown(
        f"""
        <div class="text-preview">
            <pre>{preview_text or "No extracted text available."}</pre>
        </div>
        """,
        unsafe_allow_html=True,
    )


def read_text_preview(file_path, max_characters=3000):
    """Read a short text preview from a file if possible."""
    if not file_exists(file_path):
        return ""
    return Path(file_path).read_text(encoding="utf-8", errors="ignore")[:max_characters]


def preview_docx(file_path):
    """Render a DOCX file with formatted paragraphs, runs, and tables as HTML in Streamlit."""
    try:
        import docx
        doc = docx.Document(file_path)
        html_output = []
        
        # Style sheet container
        html_output.append("<div style='background: white; color: #1e293b; padding: 2rem; border-radius: 8px; border: 1px solid #e2e8f0; font-family: system-ui, sans-serif;'>")
        
        for element in doc.element.body:
            if element.tag.endswith('p'):
                # Extract paragraph
                p = None
                for doc_p in doc.paragraphs:
                    if doc_p._p is element:
                        p = doc_p
                        break
                if p:
                    text_runs = []
                    for run in p.runs:
                        run_html = html.escape(run.text)
                        if run.bold:
                            run_html = f"<strong>{run_html}</strong>"
                        if run.italic:
                            run_html = f"<em>{run_html}</em>"
                        if run.underline:
                            run_html = f"<u>{run_html}</u>"
                        text_runs.append(run_html)
                    
                    full_text = "".join(text_runs)
                    if not full_text.strip():
                        continue
                        
                    if p.style.name.startswith("Heading"):
                        level = p.style.name.replace("Heading", "").strip()
                        level = level if level in ("1", "2", "3", "4") else "2"
                        html_output.append(f"<h{level} style='margin-top: 1.5rem; margin-bottom: 0.5rem; color: #0f172a; font-weight: 700;'>{full_text}</h{level}>")
                    else:
                        html_output.append(f"<p style='margin-bottom: 1rem; line-height: 1.6;'>{full_text}</p>")
                        
            elif element.tag.endswith('tbl'):
                # Extract table
                for doc_t in doc.tables:
                    if doc_t._tbl is element:
                        html_output.append("<table style='width: 100%; border-collapse: collapse; margin-bottom: 1.5rem;'>")
                        for row in doc_t.rows:
                            html_output.append("<tr>")
                            for cell in row.cells:
                                html_output.append("<td style='border: 1px solid #cbd5e1; padding: 0.5rem 0.75rem;'>")
                                cell_texts = []
                                for cp in cell.paragraphs:
                                    cell_runs = []
                                    for run in cp.runs:
                                        run_html = html.escape(run.text)
                                        if run.bold:
                                            run_html = f"<strong>{run_html}</strong>"
                                        if run.italic:
                                            run_html = f"<em>{run_html}</em>"
                                        cell_runs.append(run_html)
                                    cell_texts.append("".join(cell_runs))
                                html_output.append("<br>".join(cell_texts))
                                html_output.append("</td>")
                            html_output.append("</tr>")
                        html_output.append("</table>")
                        break
                        
        html_output.append("</div>")
        st.markdown("\n".join(html_output), unsafe_allow_html=True)
    except Exception as e:
        st.error(f"Failed to render Word document preview: {e}")


def preview_pptx(file_path):
    """Render a PowerPoint presentation slides as simple HTML cards in Streamlit."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        html_output = []
        
        for idx, slide in enumerate(prs.slides):
            html_output.append(f"<div style='background: white; border: 1px solid #e2e8f0; border-radius: 8px; padding: 1.5rem; margin-bottom: 1.5rem; box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.1); font-family: system-ui, sans-serif;'>")
            html_output.append(f"<div style='font-size: 0.75rem; text-transform: uppercase; letter-spacing: 0.05em; color: #64748b; margin-bottom: 0.5rem;'>Slide {idx + 1}</div>")
            
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    p_text_runs = []
                    for run in paragraph.runs:
                        run_html = html.escape(run.text)
                        if run.font.bold:
                            run_html = f"<strong>{run_html}</strong>"
                        if run.font.italic:
                            run_html = f"<em>{run_html}</em>"
                        p_text_runs.append(run_html)
                        
                    full_p_text = "".join(p_text_runs)
                    if not full_p_text.strip():
                        continue
                        
                    if paragraph.level == 0:
                        html_output.append(f"<h3 style='color: #0f172a; margin-top: 0.5rem; margin-bottom: 0.5rem;'>{full_p_text}</h3>")
                    else:
                        margin_left = paragraph.level * 20
                        html_output.append(f"<p style='margin-left: {margin_left}px; margin-bottom: 0.25rem; line-height: 1.5; color: #334155;'>• {full_p_text}</p>")
            
            html_output.append("</div>")
            
        st.markdown("\n".join(html_output), unsafe_allow_html=True)
    except Exception as e:
        st.error(f"Failed to render PowerPoint preview: {e}")


def preview_xlsx(file_path):
    """Render Excel spreadsheets as formatted HTML tables in tabs."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        sheet_names = wb.sheetnames
        
        tabs = st.tabs(sheet_names)
        for idx, tab in enumerate(tabs):
            with tab:
                sheet = wb[sheet_names[idx]]
                html_output = ["<table style='width: 100%; border-collapse: collapse; font-size: 0.875rem; text-align: left; font-family: system-ui, sans-serif;'>"]
                
                for r_idx, row in enumerate(sheet.iter_rows(max_row=100, max_col=20, values_only=True)):
                    if all(v is None for v in row):
                        continue
                    html_output.append("<tr>")
                    for val in row:
                        display_val = html.escape(str(val)) if val is not None else ""
                        if r_idx == 0:
                            html_output.append(f"<th style='background: #f8fafc; border: 1px solid #cbd5e1; padding: 0.5rem; font-weight: 700; color: #1e293b;'>{display_val}</th>")
                        else:
                            html_output.append(f"<td style='border: 1px solid #e2e8f0; padding: 0.5rem; color: #475569;'>{display_val}</td>")
                    html_output.append("</tr>")
                
                html_output.append("</table>")
                st.markdown("\n".join(html_output), unsafe_allow_html=True)
    except Exception as e:
        st.error(f"Failed to render Excel preview: {e}")
