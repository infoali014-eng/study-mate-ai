from pathlib import Path
import os

import fitz


IMAGE_TYPES = {"PNG", "JPG", "JPEG", "WEBP"}
TEXT_TYPES = {"TXT", "MD", "CSV", "JSON"}
SUPPORTED_FILE_TYPES = {"PDF", *IMAGE_TYPES, "DOCX", "PPTX", "PPT", "XLSX", *TEXT_TYPES}
MAX_PDF_OCR_PAGES = int(os.getenv("STUDYMATE_MAX_PDF_OCR_PAGES", "3"))
PDF_OCR_ZOOM = float(os.getenv("STUDYMATE_PDF_OCR_ZOOM", "1.4"))
OCR_UNAVAILABLE_MESSAGE = (
    "OCR support is not available yet on this deployment. Image-based text may not be extracted."
)


def ocr_status():
    """Return a friendly OCR availability status."""
    try:
        import pytesseract

        pytesseract.get_tesseract_version()
        return "OCR available"
    except Exception:
        return "OCR not available"


def _ocr_image(image):
    """Return OCR text from a PIL image, or a warning when OCR is unavailable."""
    try:
        import pytesseract

        text = pytesseract.image_to_string(image)
        return text.strip(), ""
    except Exception:
        return "", OCR_UNAVAILABLE_MESSAGE


def _extract_pdf(file_path):
    """Extract selectable PDF text and OCR a small sample of scanned pages.

    Long scanned PDFs can take many minutes if every page is rendered and sent to
    OCR. The app keeps uploads responsive by extracting selectable text from all
    pages, then OCRing only a few low-text pages as a helpful sample.
    """
    text_parts = []
    warnings = []
    page_count = 0
    ocr_used = False
    scanned_pages = 0
    ocr_attempts = 0
    ocr_warning_added = False

    with fitz.open(file_path) as document:
        page_count = document.page_count
        for page_number, page in enumerate(document, start=1):
            page_text = page.get_text("text").strip()
            if page_text:
                text_parts.append(f"\n--- Page {page_number} ---\n{page_text}")

            if len(page_text) < 40:
                scanned_pages += 1
                if ocr_attempts >= MAX_PDF_OCR_PAGES:
                    continue

                ocr_attempts += 1
                try:
                    from PIL import Image

                    pixmap = page.get_pixmap(
                        matrix=fitz.Matrix(PDF_OCR_ZOOM, PDF_OCR_ZOOM),
                        alpha=False,
                    )
                    image = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
                    ocr_text, warning = _ocr_image(image)
                    if ocr_text:
                        text_parts.append(f"\n--- Page {page_number} OCR ---\n{ocr_text}")
                        ocr_used = True
                    elif warning and not ocr_warning_added:
                        warnings.append(
                            "This PDF appears to be scanned/image-based. "
                            f"{warning}"
                        )
                        ocr_warning_added = True
                except Exception:
                    if not ocr_warning_added:
                        warnings.append(
                            "This PDF appears to be scanned/image-based. "
                            f"{OCR_UNAVAILABLE_MESSAGE}"
                        )
                        ocr_warning_added = True

    method = "pdf_ocr" if ocr_used else "pdf_text"
    if scanned_pages > MAX_PDF_OCR_PAGES:
        warnings.append(
            "Long or scanned PDF detected. To keep upload fast, OCR was limited to "
            f"{MAX_PDF_OCR_PAGES} page(s). Selectable text was still extracted from all pages."
        )
    if scanned_pages and not ocr_used and not warnings:
        warnings.append(
            "Some PDF pages appear scanned/image-based, but no OCR text was extracted."
        )
    return "\n".join(text_parts).strip(), page_count, method, warnings


def _extract_image(file_path):
    """Extract text from an uploaded image with OCR when available."""
    warnings = []
    try:
        from PIL import Image

        with Image.open(file_path) as image:
            text, warning = _ocr_image(image)
    except Exception:
        text, warning = "", OCR_UNAVAILABLE_MESSAGE

    if warning:
        warnings.append(
            "Image uploaded, but text extraction failed. You can still preview/download it. "
            f"{warning}"
        )

    return text.strip(), 1, "image_ocr" if text else "ocr_unavailable", warnings


def _extract_docx(file_path):
    """Extract paragraphs and table text from a DOCX file."""
    try:
        from docx import Document
    except Exception:
        return "", 0, "docx_unavailable", ["DOCX support is not available on this deployment."]

    document = Document(file_path)
    parts = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            style_name = (paragraph.style.name or "").lower() if paragraph.style else ""
            prefix = "## " if "heading" in style_name else ""
            parts.append(f"{prefix}{text}")

    for table_index, table in enumerate(document.tables, start=1):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"\nTable {table_index}:\n" + "\n".join(rows))

    return "\n\n".join(parts).strip(), 0, "docx_text", []


def _shape_text(shape):
    """Extract text from a PowerPoint shape, including tables."""
    pieces = []
    if getattr(shape, "has_text_frame", False) and shape.text_frame:
        text = shape.text_frame.text.strip()
        if text:
            pieces.append(text)

    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            pieces.append("\n".join(rows))

    return "\n".join(pieces).strip()


def _extract_pptx(file_path):
    """Extract slide text from a PPTX file."""
    try:
        from pptx import Presentation
    except Exception:
        return "", 0, "pptx_unavailable", ["PPTX support is not available on this deployment."]

    try:
        presentation = Presentation(file_path)
    except Exception as e:
        return "", 0, "pptx_error", [f"Could not read presentation: {e}. Only newer .pptx files are supported."]
    parts = []
    warnings = []
    image_only_slides = 0

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            text = _shape_text(shape)
            if text:
                slide_texts.append(text)

        if slide_texts:
            title = slide_texts[0]
            content = "\n".join(slide_texts[1:]).strip()
            slide_block = f"Slide {slide_index}:\nTitle: {title}"
            if content:
                slide_block += f"\nContent:\n{content}"
            parts.append(slide_block)
        else:
            image_only_slides += 1

    if image_only_slides:
        warnings.append("Some slides appear image-based. OCR for slide images may be limited.")

    return "\n\n".join(parts).strip(), len(presentation.slides), "pptx_text", warnings


def _extract_xlsx(file_path):
    """Extract readable cell text from an XLSX workbook."""
    try:
        from openpyxl import load_workbook
    except Exception:
        return "", 0, "xlsx_unavailable", ["XLSX support is not available on this deployment."]

    workbook = load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    sheet_count = 0

    for sheet in workbook.worksheets:
        sheet_count += 1
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if values:
                rows.append(" | ".join(values))
            if len(rows) >= 400:
                rows.append("... additional rows skipped for fast upload ...")
                break

        if rows:
            parts.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))

    workbook.close()
    return "\n\n".join(parts).strip(), sheet_count, "xlsx_text", []


def _extract_text_file(file_path):
    """Read TXT or Markdown files safely."""
    path = Path(file_path)
    try:
        return path.read_text(encoding="utf-8", errors="ignore").strip()
    except Exception:
        return path.read_text(encoding="latin-1", errors="ignore").strip()


def process_uploaded_file(file_path, file_type):
    """Process any supported study material into extracted text and metadata."""
    path = Path(file_path)
    clean_type = (file_type or path.suffix.replace(".", "") or "").upper()
    warnings = []

    if clean_type not in SUPPORTED_FILE_TYPES:
        return {
            "text": "",
            "page_count": 0,
            "method": "unsupported",
            "status": "Unsupported file type",
            "warnings": [f"{clean_type or 'This file'} is not supported."],
        }

    if clean_type == "PDF":
        text, count, method, warnings = _extract_pdf(path)
    elif clean_type in IMAGE_TYPES:
        text, count, method, warnings = _extract_image(path)
    elif clean_type == "DOCX":
        text, count, method, warnings = _extract_docx(path)
    elif clean_type in {"PPTX", "PPT"}:
        text, count, method, warnings = _extract_pptx(path)
    elif clean_type == "XLSX":
        text, count, method, warnings = _extract_xlsx(path)
    else:
        text = _extract_text_file(path)
        count = 1
        method = "plain_text"

    if text and method in {"pdf_ocr", "image_ocr"}:
        status = "OCR used"
    elif text:
        status = "Text extracted"
    elif method == "ocr_unavailable":
        status = "OCR unavailable"
    else:
        status = "Text not found"

    if not text and not warnings:
        warnings.append(
            "File uploaded, but no readable text was extracted. You can still preview it, but AI chat may not use its content."
        )

    return {
        "text": text,
        "page_count": count,
        "method": method,
        "status": status,
        "warnings": warnings,
    }
